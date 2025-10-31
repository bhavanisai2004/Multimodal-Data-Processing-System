import streamlit as st
import google.generativeai as genai
import os
from PIL import Image
import PyPDF2
import docx
from pptx import Presentation
from pytube import YouTube
import tempfile

# -----------------------------
# 🔑 Gemini API Configuration
# -----------------------------
GEMINI_API_KEY = "AIzaSyBmRtKZ5yvqVLyhsHFWcdUPyrZXcRYSWPQ"  # Replace with your Gemini API key
genai.configure(api_key=GEMINI_API_KEY)

# Use a working model (Gemini 2.5 Flash or Pro)
MODEL_NAME = "models/gemini-2.5-flash"
model = genai.GenerativeModel(MODEL_NAME)

# -----------------------------
# 🌟 Streamlit App Setup
# -----------------------------
st.set_page_config(page_title="Multimodal Data Processing System", layout="wide")
st.title("🧠 Multimodal Data Processing System")
st.write("Upload text, document, image, or YouTube link and ask questions!")

# Initialize session state
if "uploaded_docs" not in st.session_state:
    st.session_state.uploaded_docs = {}

# -----------------------------
# 📥 File Upload Section
# -----------------------------
uploaded_files = st.file_uploader(
    "Upload files (PDF, DOCX, PPTX, TXT, JPG, PNG)",
    type=["pdf", "docx", "pptx", "txt", "jpg", "jpeg", "png"],
    accept_multiple_files=True
)

if uploaded_files:
    for file in uploaded_files:
        st.session_state.uploaded_docs[file.name] = file
    st.success(f"✅ {len(uploaded_files)} files uploaded successfully!")

# -----------------------------
# 🎥 YouTube Link Input
# -----------------------------
yt_link = st.text_input("Or paste a YouTube video link (optional):")

if yt_link:
    try:
        yt = YouTube(yt_link)
        st.write(f"🎬 **Title:** {yt.title}")
        st.video(yt_link)
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".mp4")
        yt.streams.filter(only_audio=True).first().download(filename=temp_file.name)
        st.session_state.uploaded_docs[yt.title] = temp_file.name
        st.success("✅ YouTube audio added!")
    except Exception as e:
        st.error(f"❌ Error fetching YouTube: {str(e)}")

# -----------------------------
# 💬 Ask a Question
# -----------------------------
query = st.text_input("Ask a question about your uploaded data:")

def extract_text_from_pdf(file):
    reader = PyPDF2.PdfReader(file)
    text = "".join([page.extract_text() for page in reader.pages if page.extract_text()])
    return text

def extract_text_from_docx(file):
    doc = docx.Document(file)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# -----------------------------
# 🧠 Process Uploaded Content
# -----------------------------
if st.button("🔍 Analyze"):
    if not st.session_state.uploaded_docs:
        st.warning("Please upload at least one file or add a YouTube link.")
    elif not query:
        st.warning("Please enter a question to analyze.")
    else:
        combined_text = ""
        for name, file in st.session_state.uploaded_docs.items():
            st.write(f"📄 Processing: **{name}**")

            if isinstance(file, str) and file.endswith(".mp4"):
                combined_text += f"\n(Audio from {name} was processed.)"
            elif name.lower().endswith(".pdf"):
                combined_text += extract_text_from_pdf(file)
            elif name.lower().endswith(".docx"):
                combined_text += extract_text_from_docx(file)
            elif name.lower().endswith(".pptx"):
                combined_text += extract_text_from_pptx(file)
            elif name.lower().endswith((".jpg", ".jpeg", ".png")):
                img = Image.open(file)
                result = model.generate_content(["Describe this image:", img])
                combined_text += result.text
            else:
                combined_text += file.read().decode("utf-8")

        st.write("🧠 Generating response...")
        response = model.generate_content(f"{combined_text}\n\nUser question: {query}")
        st.subheader("✨ Answer:")
        st.write(response.text)

# -----------------------------
# 🧹 Clear Uploaded Files
# -----------------------------
if st.button("🧹 Clear All"):
    st.session_state.uploaded_docs = {}
    st.success("All files cleared!")
